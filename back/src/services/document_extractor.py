import os
import logging
from typing import Optional

logger = logging.getLogger(__name__)

class DocumentExtractorService:
    """Service to extract text from documents and convert to Markdown"""

    def extract_to_markdown(self, file_path: str) -> Optional[str]:
        """
        Converts a document (PDF, DOCX, XLSX, PPTX) to Markdown text.

        Args:
            file_path: Absolute path to the local file

        Returns:
            Markdown content as string, or None if extraction fails
        """
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return None

        extension = os.path.splitext(file_path)[1].lower()
        extractors = {
            ".pdf": self._extract_pdf,
            ".docx": self._extract_docx,
            ".pptx": self._extract_pptx,
            ".xlsx": self._extract_xlsx,
        }

        extractor = extractors.get(extension)
        if not extractor:
            logger.error(f"Unsupported file extension: {extension}")
            return None

        try:
            logger.info(f"Extracting content from: {file_path}")
            return extractor(file_path)
        except Exception as e:
            logger.error(f"Error converting document {file_path}: {str(e)}")
            return None

    def _extract_pdf(self, file_path: str) -> str:
        from PyPDF2 import PdfReader

        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                pages.append(text.strip())
        return "\n\n".join(pages)

    def _extract_docx(self, file_path: str) -> str:
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        document = Document(file_path)
        parts = []
        for element in document.element.body:
            if element.tag.endswith("}p"):
                paragraph = Paragraph(element, document)
                text = paragraph.text.strip()
                if not text:
                    continue
                style = (paragraph.style.name or "").lower()
                if style.startswith("heading"):
                    level = "".join(ch for ch in style if ch.isdigit()) or "1"
                    parts.append(f"{'#' * min(int(level), 6)} {text}")
                else:
                    parts.append(text)
            elif element.tag.endswith("}tbl"):
                table = Table(element, document)
                parts.append(self._table_to_markdown(
                    [[cell.text.strip() for cell in row.cells] for row in table.rows]
                ))
        return "\n\n".join(parts)

    def _extract_pptx(self, file_path: str) -> str:
        from pptx import Presentation

        presentation = Presentation(file_path)
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        texts.append(text)
                if shape.has_table:
                    texts.append(self._table_to_markdown(
                        [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    ))
            if texts:
                slides.append(f"## Slide {index}\n\n" + "\n\n".join(texts))
        return "\n\n".join(slides)

    def _extract_xlsx(self, file_path: str) -> str:
        from openpyxl import load_workbook

        workbook = load_workbook(file_path, read_only=True, data_only=True)
        sheets = []
        for worksheet in workbook.worksheets:
            rows = [
                ["" if cell is None else str(cell) for cell in row]
                for row in worksheet.iter_rows(values_only=True)
                if any(cell is not None and str(cell).strip() for cell in row)
            ]
            if rows:
                sheets.append(f"## {worksheet.title}\n\n" + self._table_to_markdown(rows))
        workbook.close()
        return "\n\n".join(sheets)

    def _table_to_markdown(self, rows: list) -> str:
        if not rows:
            return ""
        width = max(len(row) for row in rows)
        normalized = [row + [""] * (width - len(row)) for row in rows]
        cells = [[str(cell).replace("|", "\\|").replace("\n", " ") for cell in row] for row in normalized]
        lines = ["| " + " | ".join(cells[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
        lines.extend("| " + " | ".join(row) + " |" for row in cells[1:])
        return "\n".join(lines)

    def cleanup_local_file(self, file_path: str):
        """Removes the temporary local file"""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                logger.info(f"Temporary file removed: {file_path}")
        except Exception as e:
            logger.error(f"Error removing temporary file {file_path}: {str(e)}")
